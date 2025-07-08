"""
PowerPoint Service

Handles loading and extracting data from .pptx files using python-pptx.
This service contains no UI logic.
"""
from typing import List, Tuple, Optional
from pptx import Presentation
from pptx.util import Inches

class PowerPointService:
    """A service for interacting with PowerPoint presentations."""

    def load_presentation(self, file_path: str) -> Optional[Presentation]:
        """Loads a .pptx file into a Presentation object."""
        try:
            return Presentation(file_path)
        except Exception as e:
            print(f"Error loading presentation at {file_path}: {e}")
            return None

    def get_slide_data(self, prs: Presentation) -> List[Tuple[str, str]]:
        """
        Extracts the ID and title from each slide in a presentation.
        
        Returns:
            A list of tuples, where each tuple is (slide_id, slide_title).
        """
        if not prs:
            return []
        
        slide_data = []
        for slide in prs.slides:
            slide_id = str(slide.slide_id)
            # We'll use a simpler title extraction for now. This can be customized.
            slide_title = slide.shapes.title.text if slide.shapes.title else f"Slide {len(slide_data) + 1}"
            slide_data.append((slide_id, slide_title))
            
        return slide_data