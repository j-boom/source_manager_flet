"""
PowerPoint Service for the Source Manager Flet application.
Handles loading and extracting data from PowerPoint presentations.
"""

import logging
from pathlib import Path
from pptx import Presentation
from typing import List, Tuple, Optional, Dict
from pprint import pprint

class PowerPointManager:
    """
    Service for handling PowerPoint operations in the Flet application.
    """

    def __init__(self, controller):
        """Initialize the PowerPoint service."""
        self.controller = controller
        self.logger = logging.getLogger(self.__class__.__name__)

    def get_slides_from_file(self, filepath: str) -> List[Dict] | None:
        """
        Reads a .pptx file and extracts a list of slides with their IDs and titles.

        Returns:
            A list of dictionaries, e.g.,
            [{'slide_id': 256, 'title': 'Introduction', 'sources': []}, ...]
            or None if the file cannot be read.
        """
        self.logger.info(f"Reading slides from: {filepath}")
        file_path = Path(filepath)
        if not file_path.exists():
            self.logger.warning(f"File not found: {filepath}")
            return None

        try:
            prs = Presentation(file_path)
            slides_data = []
            for i, slide in enumerate(prs.slides):
                # Use the unique slide_id from python-pptx
                slide_id = slide.slide_id
                # Use the new, more robust title extraction logic
                title = self._extract_slide_title(slide, i)

                slides_data.append(
                    {
                        "slide_id": slide_id,
                        "title": title,
                        "sources": [],  # This will be populated by the controller
                    }
                )

            self.logger.info(f"Successfully extracted {len(slides_data)} slides.")
            pprint(slides_data)
            return slides_data

        except Exception as e:
            self.logger.error(
                f"Failed to read presentation file at {filepath}: {e}", exc_info=True
            )
            return None

    def _extract_slide_title(self, slide, slide_index: int) -> str:
        """
        Extract the title from a slide using multiple strategies.

        Args:
            slide: The slide object
            slide_index: Index of the slide (0-based)

        Returns:
            The slide title or a default title
        """
        # Strategy 1: Check if slide has a title shape
        if (
            hasattr(slide, "shapes")
            and hasattr(slide.shapes, "title")
            and slide.shapes.title
        ):
            if hasattr(slide.shapes.title, "text") and slide.shapes.title.text.strip():
                return slide.shapes.title.text.strip()

        # Strategy 2: Look for text at the top of the slide (title position)
        title_text = self._find_title_by_position(slide)
        if title_text:
            return title_text

        # Strategy 3: Find the largest text box
        title_text = self._find_title_by_size(slide)
        if title_text:
            return title_text

        # Fallback: Use default title
        return f"Slide {slide_index + 1}"

    def _find_title_by_position(self, slide) -> Optional[str]:
        """
        Find title by looking for text shapes in the top portion of the slide.
        """
        try:
            # Look for text shapes in the top 25% of the slide
            for shape in slide.shapes:
                if (
                    hasattr(shape, "has_text_frame")
                    and shape.has_text_frame
                    and hasattr(shape, "top")
                    and hasattr(shape, "text")
                ):
                    # Check if shape is in the top portion (top < 25% of slide height)
                    if shape.top < 914400 * 2:  # Roughly top 2 inches
                        text = shape.text.strip()
                        if text and len(text) < 100:  # Reasonable title length
                            return text
        except Exception:
            pass
        return None

    def _find_title_by_size(self, slide) -> Optional[str]:
        """
        Find title by looking for the largest text box.
        """
        try:
            largest_text = ""
            largest_area = 0

            for shape in slide.shapes:
                if (
                    hasattr(shape, "has_text_frame")
                    and shape.has_text_frame
                    and hasattr(shape, "width")
                    and hasattr(shape, "height")
                    and hasattr(shape, "text")
                ):
                    area = shape.width * shape.height
                    text = shape.text.strip()
                    if area > largest_area and text and len(text) < 100:
                        largest_area = area
                        largest_text = text

            return largest_text if largest_text else None
        except Exception:
            pass
        return None
